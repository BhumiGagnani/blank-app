import collections
import collections.abc
from pptx import Presentation
import comtypes.client
import os

file = 'C:/Users/Heena/anaconda3/envs/ppt2md/FEMAR_A3_2_1_CM_General.pptx' #Name and Location of the File

ppt = Presentation(file)

# List to hold notes
notes = []
image_dir = 'C:/Users/Heena/Desktop/Masters TUD/WHK/Slides_as_images/FEMAR_A3_2_1_CM_General' #Name and Location of the folder where the slides are saved in .jpg format

# Extract notes from each slide
for page, slide in enumerate(ppt.slides):
   
    # Get the presenter notes
    textNote = slide.notes_slide.notes_text_frame.text
    notes.append((page + 1, textNote))  # Use page + 1 for 1-based indexing

# Define the output Markdown file path
output_file = 'FEMAR_A3_2_1_CM_General6.md'

# Write the notes and images to the Markdown file
with open(output_file, 'w', encoding='utf-8') as f:
    for page, textNote in notes:
        f.write(f"---\n\n")  # Markdown header for each slide
        f.write(f"{textNote}\n\n")      # Write the notes
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            textNote = slide.notes_slide.notes_text_frame.text
        else:
            textNote = "(No notes for this slide)"
        f.write("---\n\n")               # Separator between slides

output_file_path = os.path.abspath(output_file)
print(f"Presentation with notes have been saved to {output_file_path}")
